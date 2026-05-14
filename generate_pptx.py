#!/usr/bin/env python3
"""Generate a professional PowerPoint presentation for À l'Heure newspaper."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# Brand colors
BLUE = RGBColor(14, 43, 98)
RED = RGBColor(200, 48, 42)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(245, 243, 238)
DARK_INK = RGBColor(26, 26, 26)

# Slide dimensions (widescreen 16:9)
SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)


def create_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs


def add_rectangle(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    """Add a rectangle shape to the slide."""
    from pptx.util import Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()

    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()

    return shape


def add_text_box(slide, text, left, top, width, height,
                  font_size=14, font_color=DARK_INK, bold=False, italic=False,
                  align=PP_ALIGN.LEFT, font_name="Calibri", word_wrap=True,
                  vertical_anchor=None):
    """Add a text box to the slide."""
    from pptx.enum.text import MSO_ANCHOR
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap

    if vertical_anchor:
        tf.vertical_anchor = vertical_anchor

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name

    return txBox


def add_text_to_shape(shape, text, font_size=14, font_color=DARK_INK, bold=False,
                       align=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add text to an existing shape."""
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.name = font_name


def add_paragraph_to_shape(shape, text, font_size=12, font_color=DARK_INK, bold=False,
                             align=PP_ALIGN.LEFT, font_name="Calibri", space_before=None):
    """Add a paragraph to an existing shape's text frame."""
    tf = shape.text_frame
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.name = font_name
    return p


def add_header_bar(slide, slide_num=None, show_logo=True):
    """Add the blue header bar to a content slide."""
    # Blue header bar
    header = add_rectangle(slide, 0, 0, SLIDE_WIDTH, Inches(0.8), fill_color=BLUE)

    if show_logo:
        # Logo text in top-left
        add_text_box(slide, "À L'HEURE",
                     Inches(0.2), Inches(0.1), Inches(2.5), Inches(0.6),
                     font_size=18, font_color=WHITE, bold=True,
                     font_name="Calibri")

    if slide_num:
        # Slide number in top-right
        add_text_box(slide, str(slide_num),
                     Inches(12.8), Inches(0.1), Inches(0.4), Inches(0.6),
                     font_size=12, font_color=WHITE, bold=False,
                     align=PP_ALIGN.RIGHT)

    return header


def add_slide_title(slide, title):
    """Add a styled slide title below the header."""
    # Red accent line
    add_rectangle(slide, Inches(0.4), Inches(0.95), Inches(0.08), Inches(0.55), fill_color=RED)

    # Title text
    add_text_box(slide, title,
                 Inches(0.6), Inches(0.9), Inches(12), Inches(0.65),
                 font_size=26, font_color=BLUE, bold=True,
                 font_name="Calibri")


def add_footer(slide, text="alheure.info"):
    """Add a subtle footer."""
    # Light footer bar
    add_rectangle(slide, 0, Inches(7.2), SLIDE_WIDTH, Inches(0.3), fill_color=LIGHT_GRAY)
    add_text_box(slide, text,
                 Inches(0.3), Inches(7.2), Inches(6), Inches(0.3),
                 font_size=9, font_color=RGBColor(120, 120, 120),
                 font_name="Calibri")
    add_text_box(slide, "© 2026 À l'Heure – Tous droits réservés",
                 Inches(6.5), Inches(7.2), Inches(6.5), Inches(0.3),
                 font_size=9, font_color=RGBColor(120, 120, 120),
                 align=PP_ALIGN.RIGHT, font_name="Calibri")


# ─────────────────────────────────────────────
# SLIDE 1: Cover
# ─────────────────────────────────────────────
def slide_cover(prs):
    blank_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(blank_layout)

    # Full blue background
    add_rectangle(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, fill_color=BLUE)

    # Decorative red band at top
    add_rectangle(slide, 0, 0, SLIDE_WIDTH, Inches(0.12), fill_color=RED)

    # Decorative red band at bottom
    add_rectangle(slide, 0, Inches(7.38), SLIDE_WIDTH, Inches(0.12), fill_color=RED)

    # Subtle light line accent (horizontal divider)
    add_rectangle(slide, Inches(1.2), Inches(4.2), Inches(10.9), Inches(0.025),
                  fill_color=RGBColor(255, 255, 255))

    # Small red square decoration (top-left corner area)
    add_rectangle(slide, Inches(0.5), Inches(1.5), Inches(0.2), Inches(0.2), fill_color=RED)

    # Newspaper-style column lines (decorative)
    for x_pos in [Inches(4.5), Inches(8.83)]:
        add_rectangle(slide, x_pos, Inches(1.4), Inches(0.015), Inches(5.5),
                      fill_color=RGBColor(40, 70, 130))

    # Main title "À l'Heure"
    add_text_box(slide, "À L'HEURE",
                 Inches(0.8), Inches(1.4), Inches(11.7), Inches(1.9),
                 font_size=88, font_color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER, font_name="Calibri")

    # Red underline accent under title
    add_rectangle(slide, Inches(2.5), Inches(3.15), Inches(8.3), Inches(0.07), fill_color=RED)

    # Subtitle
    add_text_box(slide, "LE QUOTIDIEN D'INFORMATION DU SÉNÉGAL ET DE L'AFRIQUE",
                 Inches(0.8), Inches(3.35), Inches(11.7), Inches(0.65),
                 font_size=16, font_color=RGBColor(200, 210, 235),
                 align=PP_ALIGN.CENTER, font_name="Calibri", bold=False)

    # Tagline in red box
    tag_box = add_rectangle(slide, Inches(2.0), Inches(4.35), Inches(9.3), Inches(0.65),
                             fill_color=RED)
    add_text_box(slide, "RIGUEUR  ·  INDÉPENDANCE  ·  PROXIMITÉ  ·  INNOVATION",
                 Inches(2.0), Inches(4.35), Inches(9.3), Inches(0.65),
                 font_size=14, font_color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER, font_name="Calibri")

    # Year
    add_text_box(slide, "2026",
                 Inches(0.8), Inches(5.3), Inches(11.7), Inches(0.55),
                 font_size=22, font_color=RGBColor(170, 185, 215),
                 align=PP_ALIGN.CENTER, font_name="Calibri")

    # Website
    add_text_box(slide, "alheure.info",
                 Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.4),
                 font_size=14, font_color=RGBColor(180, 195, 225),
                 align=PP_ALIGN.CENTER, font_name="Calibri", italic=True)


# ─────────────────────────────────────────────
# SLIDE 2: À propos
# ─────────────────────────────────────────────
def slide_apropos(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    add_rectangle(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, fill_color=WHITE)
    add_header_bar(slide, slide_num=2)
    add_slide_title(slide, "Qui sommes-nous ?")
    add_footer(slide)

    # Main content box (left side)
    content = add_rectangle(slide, Inches(0.4), Inches(1.7), Inches(7.8), Inches(5.2),
                              fill_color=LIGHT_GRAY)
    tf = content.text_frame
    tf.word_wrap = True

    from pptx.util import Pt
    from pptx.oxml.ns import qn

    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.LEFT
    r0 = p0.add_run()
    r0.text = "À l'Heure"
    r0.font.size = Pt(18)
    r0.font.bold = True
    r0.font.color.rgb = BLUE
    r0.font.name = "Calibri"

    lines = [
        ("est un quotidien sénégalais fondé en 2026, disponible "
         "en version papier et en ligne sur alheure.info.", False, 13),
        ("", False, 8),
        ("Siège social : Dakar, Sénégal", True, 12),
        ("Équipe de journalistes professionnels aguerris", False, 12),
        ("Couverture nationale et panafricaine", False, 12),
        ("", False, 8),
        ("ISSN : 3020-0001", True, 12),
        ("Prix de vente : 100 FCFA", True, 12),
    ]

    for text, is_bold, size in lines:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = is_bold
        r.font.color.rgb = DARK_INK
        r.font.name = "Calibri"

    # Padding inside shape
    from lxml import etree
    txBody = content.text_frame._txBody
    bodyPr = txBody.find(qn('a:bodyPr'))
    if bodyPr is not None:
        bodyPr.set('lIns', str(int(Inches(0.2))))
        bodyPr.set('rIns', str(int(Inches(0.2))))
        bodyPr.set('tIns', str(int(Inches(0.2))))
        bodyPr.set('bIns', str(int(Inches(0.2))))

    # Right info panel
    panel = add_rectangle(slide, Inches(8.5), Inches(1.7), Inches(4.5), Inches(5.2),
                           fill_color=BLUE)

    add_text_box(slide, "EN CHIFFRES",
                 Inches(8.7), Inches(1.9), Inches(4.1), Inches(0.5),
                 font_size=16, font_color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER, font_name="Calibri")

    # Red divider
    add_rectangle(slide, Inches(9.0), Inches(2.5), Inches(3.5), Inches(0.05), fill_color=RED)

    stats = [
        ("Fondé en", "2026"),
        ("Prix", "100 FCFA"),
        ("ISSN", "3020-0001"),
        ("Format", "Tabloïd"),
        ("Édition papier", "Lun–Ven"),
        ("Site web", "alheure.info"),
    ]

    y = Inches(2.7)
    for label, value in stats:
        add_text_box(slide, label,
                     Inches(8.7), y, Inches(2.2), Inches(0.45),
                     font_size=11, font_color=RGBColor(180, 195, 225),
                     font_name="Calibri")
        add_text_box(slide, value,
                     Inches(10.9), y, Inches(1.9), Inches(0.45),
                     font_size=13, font_color=WHITE, bold=True,
                     align=PP_ALIGN.RIGHT, font_name="Calibri")
        y += Inches(0.58)


# ─────────────────────────────────────────────
# SLIDE 3: Notre mission
# ─────────────────────────────────────────────
def slide_mission(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    add_rectangle(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, fill_color=WHITE)
    add_header_bar(slide, slide_num=3)
    add_slide_title(slide, "Notre mission")
    add_footer(slide)

    pillars = [
        ("RIGUEUR", "Journalisme de qualité,\nvérification rigoureuse\ndes faits.", BLUE),
        ("INDÉPENDANCE", "Ligne éditoriale libre\net indépendante de tout\npouvoir.", RED),
        ("PROXIMITÉ", "Ancrage sénégalais fort,\ncorrespondants dans\nles régions.", BLUE),
        ("INNOVATION", "Digital first — podcasts,\nvidéos, newsletter\nquotidienne.", RED),
    ]

    box_w = Inches(2.9)
    box_h = Inches(4.8)
    y_start = Inches(1.65)
    gap = Inches(0.38)
    x_start = Inches(0.45)

    for i, (title, desc, color) in enumerate(pillars):
        x = x_start + i * (box_w + gap)

        # Shadow effect (dark rectangle slightly offset)
        add_rectangle(slide, x + Inches(0.07), y_start + Inches(0.07),
                       box_w, box_h, fill_color=RGBColor(200, 200, 200))

        # Main box
        box = add_rectangle(slide, x, y_start, box_w, box_h, fill_color=color)

        # Number indicator
        add_text_box(slide, f"0{i+1}",
                     x + Inches(0.15), y_start + Inches(0.2), Inches(0.7), Inches(0.55),
                     font_size=28, font_color=RGBColor(255, 255, 255),
                     bold=True, font_name="Calibri",
                     align=PP_ALIGN.LEFT)

        # Pillar title
        add_text_box(slide, title,
                     x + Inches(0.15), y_start + Inches(0.85), box_w - Inches(0.3), Inches(0.65),
                     font_size=18, font_color=WHITE, bold=True,
                     font_name="Calibri")

        # Thin white line separator
        add_rectangle(slide, x + Inches(0.15), y_start + Inches(1.6),
                       box_w - Inches(0.3), Inches(0.04), fill_color=WHITE)

        # Description
        add_text_box(slide, desc,
                     x + Inches(0.15), y_start + Inches(1.8), box_w - Inches(0.3), Inches(2.7),
                     font_size=13, font_color=RGBColor(220, 225, 240),
                     font_name="Calibri", word_wrap=True)


# ─────────────────────────────────────────────
# SLIDE 4: Offre éditoriale
# ─────────────────────────────────────────────
def slide_offre(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    add_rectangle(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, fill_color=WHITE)
    add_header_bar(slide, slide_num=4)
    add_slide_title(slide, "Notre offre éditoriale")
    add_footer(slide)

    # Left panel: Rubriques
    add_rectangle(slide, Inches(0.4), Inches(1.7), Inches(5.9), Inches(5.2),
                   fill_color=LIGHT_GRAY)

    add_text_box(slide, "RUBRIQUES",
                 Inches(0.6), Inches(1.85), Inches(5.5), Inches(0.5),
                 font_size=16, font_color=BLUE, bold=True, font_name="Calibri")

    add_rectangle(slide, Inches(0.6), Inches(2.4), Inches(5.5), Inches(0.05), fill_color=RED)

    rubriques = ["Sénégal", "Afrique", "Monde", "Politique", "Économie",
                 "Société", "Sport", "Culture", "Diaspora"]

    y = Inches(2.55)
    for i, r in enumerate(rubriques):
        # bullet dot
        add_rectangle(slide, Inches(0.7), y + Inches(0.14), Inches(0.1), Inches(0.1),
                       fill_color=RED)
        add_text_box(slide, r,
                     Inches(0.95), y, Inches(5.0), Inches(0.42),
                     font_size=13, font_color=DARK_INK, font_name="Calibri")
        y += Inches(0.45)

    # Right panel: Formats
    add_rectangle(slide, Inches(6.6), Inches(1.7), Inches(6.35), Inches(5.2),
                   fill_color=BLUE)

    add_text_box(slide, "FORMATS",
                 Inches(6.8), Inches(1.85), Inches(6.0), Inches(0.5),
                 font_size=16, font_color=WHITE, bold=True, font_name="Calibri")

    add_rectangle(slide, Inches(6.8), Inches(2.4), Inches(5.9), Inches(0.05), fill_color=RED)

    formats = [
        ("Reportages", "Terrain, immersion, récits"),
        ("Enquêtes", "Investigation approfondie"),
        ("Analyses", "Décryptage et contexte"),
        ("Long format", "Dossiers thématiques"),
        ("Interviews", "Figures politiques & culturelles"),
        ("Podcasts", "Audio hebdomadaire"),
        ("Vidéos", "Format court & documentaires"),
    ]

    y = Inches(2.55)
    for fmt, desc in formats:
        add_text_box(slide, fmt,
                     Inches(6.9), y, Inches(2.5), Inches(0.42),
                     font_size=13, font_color=WHITE, bold=True, font_name="Calibri")
        add_text_box(slide, desc,
                     Inches(9.5), y, Inches(3.2), Inches(0.42),
                     font_size=11, font_color=RGBColor(180, 195, 225), font_name="Calibri")
        y += Inches(0.5)


# ─────────────────────────────────────────────
# SLIDE 5: Journal papier
# ─────────────────────────────────────────────
def slide_papier(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    add_rectangle(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, fill_color=WHITE)
    add_header_bar(slide, slide_num=5)
    add_slide_title(slide, "Le journal papier")
    add_footer(slide)

    # Newspaper mockup area (left)
    paper_bg = add_rectangle(slide, Inches(0.4), Inches(1.65), Inches(5.2), Inches(5.25),
                               fill_color=LIGHT_GRAY)

    # Fake newspaper header
    add_rectangle(slide, Inches(0.4), Inches(1.65), Inches(5.2), Inches(0.8), fill_color=BLUE)
    add_text_box(slide, "À L'HEURE",
                 Inches(0.5), Inches(1.75), Inches(5.0), Inches(0.6),
                 font_size=24, font_color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER, font_name="Calibri")

    # Fake date bar
    add_rectangle(slide, Inches(0.4), Inches(2.45), Inches(5.2), Inches(0.25), fill_color=RED)
    add_text_box(slide, "Dakar · Vendredi 15 mai 2026 · No 42 · Prix : 100 FCFA",
                 Inches(0.45), Inches(2.47), Inches(5.1), Inches(0.2),
                 font_size=7.5, font_color=WHITE, font_name="Calibri")

    # Fake headline columns
    add_text_box(slide, "MANCHETTE PRINCIPALE",
                 Inches(0.5), Inches(2.85), Inches(5.0), Inches(0.45),
                 font_size=16, font_color=DARK_INK, bold=True, font_name="Calibri")

    # Fake text blocks
    for y_pos in [Inches(3.4), Inches(3.75), Inches(4.1)]:
        add_rectangle(slide, Inches(0.5), y_pos, Inches(4.8), Inches(0.22),
                       fill_color=RGBColor(200, 200, 200))

    add_rectangle(slide, Inches(0.5), Inches(4.55), Inches(2.2), Inches(1.15),
                   fill_color=RGBColor(185, 185, 185))
    add_rectangle(slide, Inches(2.9), Inches(4.55), Inches(2.4), Inches(1.15),
                   fill_color=RGBColor(200, 200, 200))
    for y_pos in [Inches(5.8), Inches(6.1)]:
        add_rectangle(slide, Inches(2.9), y_pos, Inches(2.4), Inches(0.18),
                       fill_color=RGBColor(200, 200, 200))

    # Right side: specs
    specs = [
        ("Périodicité", "Quotidien (lundi – vendredi)"),
        ("Prix", "100 FCFA"),
        ("ISSN", "3020-0001"),
        ("Format", "Tabloïd"),
        ("Diffusion", "Kiosques de Dakar et grandes villes"),
        ("Tirage", "En développement"),
    ]

    y = Inches(1.75)
    for label, value in specs:
        # Label background
        add_rectangle(slide, Inches(6.0), y, Inches(3.2), Inches(0.55), fill_color=BLUE)
        add_text_box(slide, label,
                     Inches(6.1), y + Inches(0.08), Inches(3.0), Inches(0.4),
                     font_size=12, font_color=WHITE, bold=True, font_name="Calibri")

        # Value box
        add_rectangle(slide, Inches(9.2), y, Inches(3.8), Inches(0.55),
                       fill_color=LIGHT_GRAY, line_color=RGBColor(220, 218, 213))
        add_text_box(slide, value,
                     Inches(9.3), y + Inches(0.08), Inches(3.5), Inches(0.4),
                     font_size=12, font_color=DARK_INK, font_name="Calibri")
        y += Inches(0.7)

    # Callout quote
    add_rectangle(slide, Inches(6.0), Inches(6.15), Inches(7.0), Inches(0.65), fill_color=RED)
    add_text_box(slide, "\"L'information à l'heure, chaque matin.\"",
                 Inches(6.1), Inches(6.15), Inches(6.8), Inches(0.65),
                 font_size=14, font_color=WHITE, bold=True, italic=True,
                 align=PP_ALIGN.CENTER, font_name="Calibri")


# ─────────────────────────────────────────────
# SLIDE 6: Site web
# ─────────────────────────────────────────────
def slide_web(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    add_rectangle(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, fill_color=WHITE)
    add_header_bar(slide, slide_num=6)
    add_slide_title(slide, "Le site web — alheure.info")
    add_footer(slide)

    features = [
        ("⚡", "Articles en temps réel", "Publication instantanée dès qu'un article est validé"),
        ("🎙", "Podcasts hebdomadaires", "Émissions audio disponibles en streaming"),
        ("🎬", "Vidéos", "Formats courts et documentaires"),
        ("📖", "Mot du jour en wolof", "Valorisation de la langue nationale"),
        ("📧", "Newsletter quotidienne", "Résumé de l'actualité chaque matin"),
        ("🔒", "Administration sécurisée", "Back-office protégé pour la rédaction"),
    ]

    cols = 3
    box_w = Inches(3.8)
    box_h = Inches(1.9)
    gap_x = Inches(0.35)
    gap_y = Inches(0.3)
    x_start = Inches(0.4)
    y_start = Inches(1.75)

    for i, (icon, title, desc) in enumerate(features):
        col = i % cols
        row = i // cols
        x = x_start + col * (box_w + gap_x)
        y = y_start + row * (box_h + gap_y)

        # Feature box
        box_color = BLUE if row == 0 else LIGHT_GRAY
        text_color = WHITE if row == 0 else DARK_INK
        desc_color = RGBColor(180, 200, 235) if row == 0 else RGBColor(90, 90, 90)

        add_rectangle(slide, x, y, box_w, box_h, fill_color=box_color)

        # Icon area
        add_text_box(slide, icon,
                     x + Inches(0.15), y + Inches(0.15), Inches(0.55), Inches(0.55),
                     font_size=22, font_color=WHITE if row == 0 else BLUE,
                     font_name="Segoe UI Emoji")

        add_text_box(slide, title,
                     x + Inches(0.15), y + Inches(0.65), box_w - Inches(0.3), Inches(0.45),
                     font_size=13, font_color=text_color, bold=True, font_name="Calibri")

        add_text_box(slide, desc,
                     x + Inches(0.15), y + Inches(1.15), box_w - Inches(0.3), Inches(0.65),
                     font_size=11, font_color=desc_color, font_name="Calibri")

    # Tech stack badge
    add_rectangle(slide, Inches(0.4), Inches(6.75), Inches(12.53), Inches(0.35),
                   fill_color=BLUE)
    add_text_box(slide, "Construit avec Next.js · Tailwind CSS · API REST · Hébergement cloud",
                 Inches(0.6), Inches(6.77), Inches(12.0), Inches(0.3),
                 font_size=10, font_color=RGBColor(180, 200, 235),
                 align=PP_ALIGN.CENTER, font_name="Calibri")


# ─────────────────────────────────────────────
# SLIDE 7: Chiffres
# ─────────────────────────────────────────────
def slide_chiffres(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    add_rectangle(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, fill_color=WHITE)
    add_header_bar(slide, slide_num=7)
    add_slide_title(slide, "Nos chiffres")
    add_footer(slide)

    stats = [
        ("29+", "Articles\npubliés"),
        ("9", "Rubriques\nthématiques"),
        ("6", "Épisodes\nde podcast"),
        ("∞", "Abonnés\nen croissance"),
    ]

    box_w = Inches(2.7)
    box_h = Inches(2.8)
    gap = Inches(0.5)
    x_start = Inches(0.75)
    y_stat = Inches(1.9)

    for i, (num, label) in enumerate(stats):
        x = x_start + i * (box_w + gap)

        # Outer blue box
        add_rectangle(slide, x, y_stat, box_w, box_h, fill_color=BLUE)

        # Red top accent
        add_rectangle(slide, x, y_stat, box_w, Inches(0.12), fill_color=RED)

        # Number
        add_text_box(slide, num,
                     x + Inches(0.1), y_stat + Inches(0.3), box_w - Inches(0.2), Inches(1.3),
                     font_size=52, font_color=WHITE, bold=True,
                     align=PP_ALIGN.CENTER, font_name="Calibri")

        # Label
        add_text_box(slide, label,
                     x + Inches(0.1), y_stat + Inches(1.7), box_w - Inches(0.2), Inches(0.9),
                     font_size=13, font_color=RGBColor(180, 200, 235),
                     align=PP_ALIGN.CENTER, font_name="Calibri")

    # Domain line
    add_rectangle(slide, Inches(0.4), Inches(5.1), Inches(12.53), Inches(0.9), fill_color=LIGHT_GRAY)
    add_text_box(slide, "Domaine  ·  alheure.info",
                 Inches(0.6), Inches(5.2), Inches(6), Inches(0.65),
                 font_size=20, font_color=BLUE, bold=True, font_name="Calibri")

    # Note box
    add_rectangle(slide, Inches(0.4), Inches(6.1), Inches(12.53), Inches(0.75),
                   fill_color=LIGHT_GRAY)
    add_text_box(slide,
                 "* Chiffres au lancement — en constante progression depuis la mise en ligne.",
                 Inches(0.6), Inches(6.2), Inches(12.0), Inches(0.55),
                 font_size=11, font_color=RGBColor(100, 100, 100),
                 italic=True, font_name="Calibri")


# ─────────────────────────────────────────────
# SLIDE 8: L'équipe
# ─────────────────────────────────────────────
def slide_equipe(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    add_rectangle(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, fill_color=WHITE)
    add_header_bar(slide, slide_num=8)
    add_slide_title(slide, "L'équipe")
    add_footer(slide)

    journalists = [
        ("AÏ", "Aïssatou Ndoye", "Environnement"),
        ("MD", "Mamadou Diallo", "Politique"),
        ("FD", "Fatou Diallo", "Économie"),
        ("IS", "Ibrahima Sow", "Sport"),
        ("RM", "Rokhaya Mbaye", "Culture"),
    ]

    card_w = Inches(2.2)
    card_h = Inches(3.5)
    gap = Inches(0.25)
    x_start = Inches(0.45)
    y_start = Inches(1.75)

    for i, (initials, name, beat) in enumerate(journalists):
        x = x_start + i * (card_w + gap)

        # Card background
        add_rectangle(slide, x, y_start, card_w, card_h, fill_color=LIGHT_GRAY)

        # Avatar circle (simulated with a square)
        avatar_color = BLUE if i % 2 == 0 else RED
        add_rectangle(slide, x + Inches(0.55), y_start + Inches(0.25),
                       Inches(1.1), Inches(1.1), fill_color=avatar_color)

        # Initials
        add_text_box(slide, initials,
                     x + Inches(0.55), y_start + Inches(0.35),
                     Inches(1.1), Inches(0.85),
                     font_size=24, font_color=WHITE, bold=True,
                     align=PP_ALIGN.CENTER, font_name="Calibri")

        # Name
        add_text_box(slide, name,
                     x + Inches(0.1), y_start + Inches(1.55), card_w - Inches(0.2), Inches(0.7),
                     font_size=12, font_color=DARK_INK, bold=True,
                     align=PP_ALIGN.CENTER, font_name="Calibri")

        # Beat tag
        add_rectangle(slide, x + Inches(0.2), y_start + Inches(2.4),
                       card_w - Inches(0.4), Inches(0.4), fill_color=avatar_color)
        add_text_box(slide, beat,
                     x + Inches(0.2), y_start + Inches(2.42), card_w - Inches(0.4), Inches(0.38),
                     font_size=11, font_color=WHITE, bold=True,
                     align=PP_ALIGN.CENTER, font_name="Calibri")

    # Correspondents note
    add_rectangle(slide, Inches(11.15), Inches(1.75), Inches(1.8), Inches(3.5), fill_color=BLUE)
    add_text_box(slide, "+",
                 Inches(11.15), Inches(2.35), Inches(1.8), Inches(0.7),
                 font_size=36, font_color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER, font_name="Calibri")
    add_text_box(slide, "Correspondants\nen Afrique\net diaspora",
                 Inches(11.2), Inches(3.1), Inches(1.7), Inches(1.5),
                 font_size=11, font_color=RGBColor(180, 200, 235),
                 align=PP_ALIGN.CENTER, font_name="Calibri")

    # Team intro
    add_rectangle(slide, Inches(0.4), Inches(5.5), Inches(12.53), Inches(1.3),
                   fill_color=LIGHT_GRAY)
    add_text_box(slide,
                 "Une rédaction composée de journalistes professionnels expérimentés, "
                 "animés par la passion de l'information de qualité et l'engagement "
                 "envers les lecteurs sénégalais et africains.",
                 Inches(0.6), Inches(5.6), Inches(12.1), Inches(1.1),
                 font_size=13, font_color=DARK_INK, font_name="Calibri",
                 align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
# SLIDE 9: Nos valeurs
# ─────────────────────────────────────────────
def slide_valeurs(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    add_rectangle(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, fill_color=WHITE)
    add_header_bar(slide, slide_num=9)
    add_slide_title(slide, "Nos valeurs")
    add_footer(slide)

    values = [
        ("Journalisme d'intérêt public",
         "Nous plaçons l'intérêt général au cœur de chaque décision éditoriale."),
        ("Respect de la déontologie",
         "Notre charte éthique guide chaque journaliste dans sa pratique quotidienne."),
        ("Pluralisme et équilibre",
         "Toutes les voix méritent d'être entendues — nous garantissons la diversité des points de vue."),
        ("Ancrage sénégalais",
         "Profondément enracinés dans les réalités sociales, culturelles et économiques du Sénégal."),
        ("Ouverture africaine et mondiale",
         "Le Sénégal dans son contexte continental et international — une perspective globale."),
    ]

    y = Inches(1.75)
    for i, (title, desc) in enumerate(values):
        # Number badge
        badge_color = BLUE if i % 2 == 0 else RED
        add_rectangle(slide, Inches(0.4), y, Inches(0.5), Inches(0.75), fill_color=badge_color)
        add_text_box(slide, str(i + 1),
                     Inches(0.4), y + Inches(0.12), Inches(0.5), Inches(0.5),
                     font_size=18, font_color=WHITE, bold=True,
                     align=PP_ALIGN.CENTER, font_name="Calibri")

        # Value row background
        row_color = LIGHT_GRAY if i % 2 == 0 else WHITE
        add_rectangle(slide, Inches(0.9), y, Inches(12.0), Inches(0.75), fill_color=row_color)

        # Title
        add_text_box(slide, title,
                     Inches(1.05), y + Inches(0.04), Inches(4.0), Inches(0.38),
                     font_size=13, font_color=badge_color, bold=True, font_name="Calibri")

        # Description
        add_text_box(slide, desc,
                     Inches(5.2), y + Inches(0.12), Inches(7.5), Inches(0.55),
                     font_size=11, font_color=DARK_INK, font_name="Calibri")

        y += Inches(0.88)


# ─────────────────────────────────────────────
# SLIDE 10: Contact
# ─────────────────────────────────────────────
def slide_contact(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Blue background
    add_rectangle(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, fill_color=BLUE)

    # Red accent lines
    add_rectangle(slide, 0, 0, SLIDE_WIDTH, Inches(0.12), fill_color=RED)
    add_rectangle(slide, 0, Inches(7.38), SLIDE_WIDTH, Inches(0.12), fill_color=RED)

    # Header bar (slightly lighter blue)
    add_rectangle(slide, 0, 0, SLIDE_WIDTH, Inches(0.8),
                   fill_color=RGBColor(10, 32, 78))

    # Logo in header
    add_text_box(slide, "À L'HEURE",
                 Inches(0.2), Inches(0.1), Inches(3.0), Inches(0.6),
                 font_size=18, font_color=WHITE, bold=True, font_name="Calibri")

    # Slide number
    add_text_box(slide, "10",
                 Inches(12.8), Inches(0.1), Inches(0.4), Inches(0.6),
                 font_size=12, font_color=WHITE, bold=False,
                 align=PP_ALIGN.RIGHT)

    # Main title
    add_text_box(slide, "Contact & Réseaux",
                 Inches(0.6), Inches(0.95), Inches(12.0), Inches(0.7),
                 font_size=30, font_color=WHITE, bold=True,
                 font_name="Calibri")

    add_rectangle(slide, Inches(0.6), Inches(1.7), Inches(12.0), Inches(0.05), fill_color=RED)

    # Left column: Contact info
    contacts = [
        ("🌐  Web", "alheure.info"),
        ("✉  Email", "redaction@alheure.info"),
        ("📍  Adresse", "Dakar, Sénégal"),
    ]

    y = Inches(2.0)
    for label, value in contacts:
        add_text_box(slide, label,
                     Inches(0.6), y, Inches(3.0), Inches(0.45),
                     font_size=13, font_color=RGBColor(180, 200, 235),
                     font_name="Calibri", bold=True)
        add_text_box(slide, value,
                     Inches(3.8), y, Inches(4.0), Inches(0.45),
                     font_size=13, font_color=WHITE,
                     font_name="Calibri")
        y += Inches(0.7)

    # Right column: Social networks
    add_text_box(slide, "RÉSEAUX SOCIAUX",
                 Inches(8.0), Inches(2.0), Inches(4.8), Inches(0.5),
                 font_size=14, font_color=WHITE, bold=True,
                 font_name="Calibri")

    add_rectangle(slide, Inches(8.0), Inches(2.55), Inches(4.8), Inches(0.05), fill_color=RED)

    socials = [
        ("Facebook", "À l'Heure Sénégal"),
        ("X (Twitter)", "@alheure_sn"),
        ("Instagram", "@alheure.info"),
        ("YouTube", "À l'Heure TV"),
    ]

    y = Inches(2.75)
    for platform, handle in socials:
        add_rectangle(slide, Inches(8.0), y, Inches(2.1), Inches(0.5), fill_color=RED)
        add_text_box(slide, platform,
                     Inches(8.05), y + Inches(0.07), Inches(2.0), Inches(0.38),
                     font_size=12, font_color=WHITE, bold=True,
                     font_name="Calibri", align=PP_ALIGN.CENTER)
        add_text_box(slide, handle,
                     Inches(10.2), y + Inches(0.1), Inches(2.5), Inches(0.38),
                     font_size=12, font_color=RGBColor(200, 215, 240),
                     font_name="Calibri")
        y += Inches(0.65)

    # Closing tagline
    add_rectangle(slide, Inches(1.5), Inches(5.7), Inches(10.3), Inches(0.07), fill_color=RED)

    add_text_box(slide, "RIGUEUR  ·  INDÉPENDANCE  ·  PROXIMITÉ  ·  INNOVATION",
                 Inches(0.6), Inches(5.95), Inches(12.0), Inches(0.55),
                 font_size=16, font_color=RGBColor(180, 200, 235),
                 bold=True, align=PP_ALIGN.CENTER, font_name="Calibri")

    add_text_box(slide, "alheure.info  —  Dakar, Sénégal  —  2026",
                 Inches(0.6), Inches(6.55), Inches(12.0), Inches(0.4),
                 font_size=12, font_color=RGBColor(140, 160, 200),
                 align=PP_ALIGN.CENTER, font_name="Calibri", italic=True)


# ─────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────
def main():
    prs = create_presentation()

    slide_cover(prs)
    slide_apropos(prs)
    slide_mission(prs)
    slide_offre(prs)
    slide_papier(prs)
    slide_web(prs)
    slide_chiffres(prs)
    slide_equipe(prs)
    slide_valeurs(prs)
    slide_contact(prs)

    output_path = "/home/user/alheure/public/alheure-presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
